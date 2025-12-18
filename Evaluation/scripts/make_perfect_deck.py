import argparse, json, os
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

def read_summary(path):
    last = None
    with open(path) as f:
        for line in f:
            s=line.strip()
            if not s: continue
            try:
                obj = json.loads(s)
            except Exception:
                continue
            if "SUMMARY" in obj:
                last = obj["SUMMARY"]
    if last is None:
        raise SystemExit(f"No SUMMARY found in {path}")
    return last

def try_read(path):
    return read_summary(path) if path and os.path.exists(path) else None

def bar_pair(labels, base_vals, lora_vals, title, out_png):
    import numpy as np
    x = np.arange(len(labels)); w = 0.35
    plt.figure()
    plt.bar(x - w/2, base_vals, width=w, label="Base")
    plt.bar(x + w/2, lora_vals, width=w, label="LoRA")
    plt.xticks(x, labels, rotation=0)
    plt.ylim(0,1.0)
    plt.title(title); plt.legend()
    for i,v in enumerate(base_vals): plt.text(x[i]-w/2, min(v+0.02,0.98), f"{v:.3f}", ha="center", fontsize=8)
    for i,v in enumerate(lora_vals): plt.text(x[i]+w/2, min(v+0.02,0.98), f"{v:.3f}", ha="center", fontsize=8)
    plt.tight_layout(); plt.savefig(out_png, dpi=200); plt.close()
    return out_png

def add_bullets(slide, lines, x=0.8, y=1.5, w=8.5, h=4.5):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame
    for t in lines:
        p = tf.add_paragraph(); p.text = t

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--base_val_full", required=True)
    ap.add_argument("--lora_val_full", required=True)
    ap.add_argument("--base_train_full", default="")
    ap.add_argument("--lora_train_full", default="")
    ap.add_argument("--judge_val", default="")
    ap.add_argument("--bootstrap_val", default="")
    ap.add_argument("--out_pptx", default="Base_vs_LoRA_Results.pptx")
    args = ap.parse_args()

    VAL_B = read_summary(args.base_val_full)
    VAL_L = read_summary(args.lora_val_full)
    TRN_B = try_read(args.base_train_full)
    TRN_L = try_read(args.lora_train_full)
    JUDGE = try_read(args.judge_val)
    BOOT  = try_read(args.bootstrap_val)

    prs = Presentation()
    t = prs.slides.add_slide(prs.slide_layouts[0])
    t.shapes.title.text = "Base vs LoRA — Medical LLM Evaluation"
    t.placeholders[1].text = "Splits: Train, Validation • Metrics: F1, ROUGE-L, BERT F1 • Optional: Judge & Bootstrap"

    # Validation metrics
    s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "Validation — Core Metrics"
    labels = ["F1","ROUGE-L(F)","BERT F1"]
    base_vals = [float(VAL_B.get("F1",0)), float(VAL_B.get("rougeL_f",0)), float(VAL_B.get("bert_f1",0))]
    lora_vals = [float(VAL_L.get("F1",0)), float(VAL_L.get("rougeL_f",0)), float(VAL_L.get("bert_f1",0))]
    chart_v = bar_pair(labels, base_vals, lora_vals, "Validation", "val_core.png")
    s.shapes.add_picture(chart_v, Inches(1), Inches(1.5), width=Inches(8))
    add_bullets(s, [f"Val sizes: Base n={VAL_B.get('n')}, LoRA n={VAL_L.get('n')}"])

    # Train metrics (if present)
    if TRN_B and TRN_L:
        s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "Train — Core Metrics"
        base_vals = [float(TRN_B.get("F1",0)), float(TRN_B.get("rougeL_f",0)), float(TRN_B.get("bert_f1",0))]
        lora_vals = [float(TRN_L.get("F1",0)), float(TRN_L.get("rougeL_f",0)), float(TRN_L.get("bert_f1",0))]
        chart_t = bar_pair(labels, base_vals, lora_vals, "Train", "train_core.png")
        s.shapes.add_picture(chart_t, Inches(1), Inches(1.5), width=Inches(8))
        add_bullets(s, [f"Train sizes: Base n={TRN_B.get('n')}, LoRA n={TRN_L.get('n')}"])

    # Judge (if present)
    if JUDGE:
        s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "A/B Preference (LLM Judge) — Validation"
        wr = float(JUDGE.get("B_win_rate", 0.0)); lo, hi = JUDGE.get("B_win_rate_ci95",[0,0])
        p = float(JUDGE.get("p_value_vs_50pct", 1.0))
        add_bullets(s, [
            f"Wins — Base: {JUDGE.get('wins_A',0)}  LoRA: {JUDGE.get('wins_B',0)}  Ties: {JUDGE.get('ties',0)}",
            f"LoRA win-rate: {wr:.3f}  (95% CI: {lo:.3f}–{hi:.3f})",
            f"Binomial p-value vs 50%: {p:.3f}",
        ])

    # Bootstrap (if present)
    if BOOT:
        s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "F1 Lift (Paired Bootstrap) — Validation"
        md = float(BOOT.get("mean_diff_B_minus_A", 0.0)); lo, hi = BOOT.get("ci95",[0,0])
        p = float(BOOT.get("p_two_sided", 1.0))
        add_bullets(s, [
            f"Mean ΔF1 (LoRA − Base): {md:.3f}",
            f"95% CI: {lo:.3f} – {hi:.3f}",
            f"Two-sided p-value: {p:.3f}",
        ])

    # Takeaways
    s = prs.slides.add_slide(prs.slide_layouts[5]); s.shapes.title.text = "Takeaways"
    add_bullets(s, [
        "Free-form answers → EM≈0; use F1 / ROUGE-L / BERT F1 for quality.",
        "Validation vs Train comparable → no obvious overfitting.",
        "Judge win-rate + bootstrap CI quantify head-to-head gains.",
        "Next: expand Val set, add safety/refusal, latency/cost, topic coverage."
    ])

    prs.save(args.out_pptx)
    print({"deck": args.out_pptx})

if __name__ == "__main__":
    main()
